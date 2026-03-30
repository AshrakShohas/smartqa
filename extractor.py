"""
extractor.py
Extracts text from PDF, DOCX, PPTX, TXT files.
Supports page range selection (page_from, page_to).
Also exposes get_page_count() so the UI can show page counts before extraction.
"""

import io
import fitz          # PyMuPDF  →  pip install pymupdf
import docx          # python-docx
from pptx import Presentation
import re


# ── Public: get page count without full extraction ───────────────

def get_page_count(uploaded_file) -> int:
    """
    Returns the number of pages/slides in a file.
    Does NOT extract text — just counts pages quickly.
    """
    name = uploaded_file.name.lower()
    data = uploaded_file.read()
    uploaded_file.seek(0)   # reset for later full read

    try:
        if name.endswith(".pdf"):
            doc = fitz.open(stream=data, filetype="pdf")
            count = doc.page_count
            doc.close()
            return count

        elif name.endswith(".docx"):
            # DOCX doesn't have real pages — count paragraphs / 25 as proxy
            d = docx.Document(io.BytesIO(data))
            paras = [p for p in d.paragraphs if p.text.strip()]
            return max(1, len(paras) // 25 + 1)

        elif name.endswith((".pptx", ".ppt")):
            prs = Presentation(io.BytesIO(data))
            return len(prs.slides)

        elif name.endswith(".txt"):
            text = _decode(data)
            lines = [l for l in text.splitlines() if l.strip()]
            return max(1, len(lines) // 40 + 1)

    except Exception:
        pass

    return 1


# ── Public: extract text with page range ────────────────────────

def extract_text_from_file(uploaded_file, page_from: int = 1, page_to: int = 9999) -> tuple:
    """
    Extracts text from the file, limited to [page_from, page_to] (1-indexed, inclusive).
    Returns (text: str, metadata: dict).
    """
    name = uploaded_file.name.lower()
    data = uploaded_file.read()
    uploaded_file.seek(0)

    if name.endswith(".pdf"):
        return _extract_pdf(data, uploaded_file.name, page_from, page_to)
    elif name.endswith(".docx"):
        return _extract_docx(data, uploaded_file.name, page_from, page_to)
    elif name.endswith((".pptx", ".ppt")):
        return _extract_pptx(data, uploaded_file.name, page_from, page_to)
    elif name.endswith(".txt"):
        return _extract_txt(data, uploaded_file.name, page_from, page_to)
    else:
        return "", {"source_file": uploaded_file.name, "pages": [], "formulas": [], "sections": []}


# ── PDF ──────────────────────────────────────────────────────────

def _extract_pdf(data: bytes, filename: str, page_from: int, page_to: int) -> tuple:
    text_parts = []
    metadata = {
        "source_file": filename,
        "pages": [],
        "formulas": [],
        "sections": [],
        "total_pages": 0,
    }

    try:
        doc = fitz.open(stream=data, filetype="pdf")
        metadata["total_pages"] = doc.page_count

        # Clamp range
        p_start = max(0, page_from - 1)          # fitz is 0-indexed
        p_end   = min(doc.page_count - 1, page_to - 1)

        for page_idx in range(p_start, p_end + 1):
            page     = doc[page_idx]
            page_num = page_idx + 1
            text     = page.get_text()

            if text.strip():
                text_parts.append(f"[Page {page_num}]\n{text}")
                metadata["pages"].append(page_num)
                metadata["formulas"].extend(_detect_formulas(text))
                metadata["sections"].extend(_detect_sections(text))

        doc.close()

        # Deduplicate
        metadata["formulas"] = list(dict.fromkeys(metadata["formulas"]))
        metadata["sections"] = list(dict.fromkeys(metadata["sections"]))

    except Exception as e:
        text_parts.append(f"[PDF error: {e}]")

    return "\n".join(text_parts), metadata


# ── DOCX ─────────────────────────────────────────────────────────

def _extract_docx(data: bytes, filename: str, page_from: int, page_to: int) -> tuple:
    text_parts = []
    metadata = {
        "source_file": filename,
        "pages": [],
        "formulas": [],
        "sections": [],
        "total_pages": 0,
    }

    try:
        d = docx.Document(io.BytesIO(data))
        all_paras = [p.text for p in d.paragraphs if p.text.strip()]

        # Treat every 25 paragraphs as one "page"
        paras_per_page = 25
        total_pages    = max(1, len(all_paras) // paras_per_page + 1)
        metadata["total_pages"] = total_pages

        p_start = max(0, page_from - 1)
        p_end   = min(total_pages - 1, page_to - 1)

        selected_paras = all_paras[p_start * paras_per_page : (p_end + 1) * paras_per_page]

        for para in selected_paras:
            text_parts.append(para)
            metadata["formulas"].extend(_detect_formulas(para))
            metadata["sections"].extend(_detect_sections(para))

        # Tables
        for table in d.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    text_parts.append(row_text)

        metadata["pages"] = list(range(page_from, min(page_to, total_pages) + 1))
        metadata["formulas"] = list(dict.fromkeys(metadata["formulas"]))
        metadata["sections"] = list(dict.fromkeys(metadata["sections"]))

    except Exception as e:
        text_parts.append(f"[DOCX error: {e}]")

    return "\n".join(text_parts), metadata


# ── PPTX ─────────────────────────────────────────────────────────

def _extract_pptx(data: bytes, filename: str, page_from: int, page_to: int) -> tuple:
    text_parts = []
    metadata = {
        "source_file": filename,
        "pages": [],
        "formulas": [],
        "sections": [],
        "total_pages": 0,
    }

    try:
        prs = Presentation(io.BytesIO(data))
        metadata["total_pages"] = len(prs.slides)

        p_start = max(0, page_from - 1)
        p_end   = min(len(prs.slides) - 1, page_to - 1)

        for slide_idx in range(p_start, p_end + 1):
            slide     = prs.slides[slide_idx]
            slide_num = slide_idx + 1
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide_texts:
                content = "\n".join(slide_texts)
                text_parts.append(f"[Slide {slide_num}]\n{content}")
                metadata["pages"].append(slide_num)
                metadata["formulas"].extend(_detect_formulas(content))
                metadata["sections"].extend(_detect_sections(content))

        metadata["formulas"] = list(dict.fromkeys(metadata["formulas"]))
        metadata["sections"] = list(dict.fromkeys(metadata["sections"]))

    except Exception as e:
        text_parts.append(f"[PPTX error: {e}]")

    return "\n".join(text_parts), metadata


# ── TXT ──────────────────────────────────────────────────────────

def _extract_txt(data: bytes, filename: str, page_from: int, page_to: int) -> tuple:
    metadata = {
        "source_file": filename,
        "pages": [],
        "formulas": [],
        "sections": [],
        "total_pages": 0,
    }

    try:
        text = _decode(data)
        lines = text.splitlines()
        lines_per_page = 40
        total_pages = max(1, len(lines) // lines_per_page + 1)
        metadata["total_pages"] = total_pages

        p_start = max(0, page_from - 1)
        p_end   = min(total_pages - 1, page_to - 1)

        selected_lines = lines[p_start * lines_per_page : (p_end + 1) * lines_per_page]
        text_out = "\n".join(selected_lines)

        metadata["pages"] = list(range(page_from, min(page_to, total_pages) + 1))
        metadata["formulas"] = _detect_formulas(text_out)
        metadata["sections"] = _detect_sections(text_out)

        return text_out, metadata

    except Exception as e:
        return f"[TXT error: {e}]", metadata


# ── Helpers ──────────────────────────────────────────────────────

def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="replace")


def _detect_formulas(text: str) -> list:
    """Find mathematical formulas and equations in text."""
    patterns = [
        r'[A-Za-zΔΩαβγθπ]\s*=\s*[A-Za-z0-9+\-*/^().²³√Δ ]+',  # General equations
        r'\bE\s*=\s*mc[²2]?\b',
        r'\bF\s*=\s*ma\b',
        r'\bPV\s*=\s*nRT\b',
        r'\bΔG\s*=\s*ΔH\s*[-−]\s*TΔS\b',
        r'\bΔH\b',
        r'[∑∫∏√π∞±≤≥≠]',
        r'\d+\s*[+\-*/]\s*\d+\s*=\s*\d+',
    ]
    found = []
    for pat in patterns:
        found.extend(re.findall(pat, text))
    return [f.strip() for f in found if len(f.strip()) > 2]


def _detect_sections(text: str) -> list:
    """Detect section/chapter headers."""
    patterns = [
        r'^(Chapter|Section|Unit|Part|§)\s+[\dIVXivx]+',
        r'^\d+\.\d+\s+[A-Z]',
        r'^[A-Z][A-Za-z ]{3,40}$',
    ]
    sections = []
    for line in text.splitlines()[:30]:
        line = line.strip()
        for pat in patterns:
            if re.match(pat, line):
                sections.append(line)
    return sections
